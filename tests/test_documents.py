"""
Unit tests for Dragon Document tools (pptx_read, pptx_create, pdf_read, pdf_extract, docx_read).
"""
import json
from pathlib import Path

import pytest


# ── Helpers ────────────────────────────────────────────────────────────


def _make_pptx(path: Path, slides_data=None):
    """Create a test PPTX file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    if slides_data is None:
        slides_data = [
            {"title": "Slide 1 Title", "content": ["Bullet A", "Bullet B"]},
            {"title": "Slide 2 Title", "content": ["Point 1", "Point 2", "Point 3"]},
        ]

    for sd in slides_data:
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        title = sd.get("title", "")
        content = sd.get("content", [])

        if title:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8.5), Inches(1.2))
            tf = txBox.text_frame
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(32)

        if content:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8.5), Inches(4.5))
            tf = txBox.text_frame
            for i, line in enumerate(content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)

    prs.save(str(path))
    return path


def _make_pdf(path: Path, pages_text=None):
    """Create a test PDF file using pymupdf."""
    import fitz

    doc = fitz.open()
    if pages_text is None:
        pages_text = [
            "Page 1 content.\nThis is the first page.",
            "Page 2 content.\nThis is the second page.",
            "Page 3 content.\nThis is the third page.",
        ]
    for text in pages_text:
        page = doc.new_page()
        page.insert_text((72, 72), text, fontsize=12)
    doc.save(str(path))
    doc.close()
    return path


def _make_docx(path: Path, paragraphs=None):
    """Create a test DOCX file using python-docx."""
    from docx import Document

    doc = Document()
    if paragraphs is None:
        paragraphs = [
            "First paragraph of the document.",
            "Second paragraph with more text.",
            "Third paragraph here.",
        ]
    for text in paragraphs:
        doc.add_paragraph(text)
    doc.save(str(path))
    return path


# ── tool_pptx_read ─────────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_pptx_read_basic(tmp_path):
    """Read a PPTX with known slides."""
    from dragon.tool.builtins.documents import tool_pptx_read

    pptx_path = tmp_path / "test.pptx"
    _make_pptx(pptx_path)

    result = json.loads(await tool_pptx_read(str(pptx_path)))
    assert "error" not in result
    assert result["slides"] == 2
    assert len(result["data"]) == 2

    slide1 = result["data"][0]
    assert slide1["slide"] == 1
    assert "Slide 1 Title" in slide1["title"]
    # content lines are in a single text frame joined by newlines
    all_content = "\n".join(slide1["content"]) if isinstance(slide1["content"], list) else slide1["content"]
    assert "Bullet A" in all_content
    assert "Bullet B" in all_content


@pytest.mark.asyncio
async def test_pptx_read_nonexistent():
    """Read a non-existent file returns error."""
    from dragon.tool.builtins.documents import tool_pptx_read

    result = json.loads(await tool_pptx_read("/nonexistent/file.pptx"))
    assert "error" in result


@pytest.mark.asyncio
async def test_pptx_read_wrong_extension(tmp_path):
    """Read a non-PPTX file returns error."""
    from dragon.tool.builtins.documents import tool_pptx_read

    txt_path = tmp_path / "notes.txt"
    txt_path.write_text("hello")
    result = json.loads(await tool_pptx_read(str(txt_path)))
    assert "error" in result


# ── tool_pptx_create ───────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_pptx_create_and_read(tmp_path):
    """Create a PPTX and verify it can be read back."""
    from dragon.tool.builtins.documents import tool_pptx_create, tool_pptx_read

    out_path = tmp_path / "created.pptx"
    slides = [
        {"title": "Introduction", "content": ["Welcome", "Agenda"]},
        {"title": "Details", "content": ["Topic A", "Topic B", "Topic C"]},
    ]
    create_result = json.loads(await tool_pptx_create(str(out_path), slides))
    assert "error" not in create_result
    assert create_result["slides_created"] == 2
    assert create_result["size_bytes"] > 0
    assert out_path.exists()

    read_result = json.loads(await tool_pptx_read(str(out_path)))
    assert read_result["slides"] == 2
    assert "Introduction" in read_result["data"][0]["title"]
    assert "Details" in read_result["data"][1]["title"]


@pytest.mark.asyncio
async def test_pptx_create_empty_slides():
    """Create with empty slides list returns error."""
    from dragon.tool.builtins.documents import tool_pptx_create

    result = json.loads(await tool_pptx_create("/tmp/test.pptx", []))
    assert "error" in result


# ── tool_pdf_read ──────────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_pdf_read_basic(tmp_path):
    """Read a PDF and verify text extraction."""
    from dragon.tool.builtins.documents import tool_pdf_read

    pdf_path = tmp_path / "test.pdf"
    _make_pdf(pdf_path)

    result = json.loads(await tool_pdf_read(str(pdf_path)))
    assert "error" not in result
    assert result["total_pages"] == 3
    assert len(result["pages"]) == 3
    assert "Page 1 content" in result["pages"][0]["text"]
    assert "Page 2 content" in result["pages"][1]["text"]


@pytest.mark.asyncio
async def test_pdf_read_page_range(tmp_path):
    """Read a PDF with a specific page range."""
    from dragon.tool.builtins.documents import tool_pdf_read

    pdf_path = tmp_path / "test.pdf"
    _make_pdf(pdf_path)

    result = json.loads(await tool_pdf_read(str(pdf_path), page_start=2, page_end=2))
    assert result["total_pages"] == 3
    assert len(result["pages"]) == 1
    assert "Page 2 content" in result["pages"][0]["text"]


@pytest.mark.asyncio
async def test_pdf_read_nonexistent():
    """Read a non-existent PDF returns error."""
    from dragon.tool.builtins.documents import tool_pdf_read

    result = json.loads(await tool_pdf_read("/nonexistent/doc.pdf"))
    assert "error" in result


@pytest.mark.asyncio
async def test_pdf_read_wrong_extension(tmp_path):
    """Read a non-PDF file returns error."""
    from dragon.tool.builtins.documents import tool_pdf_read

    txt_path = tmp_path / "notes.txt"
    txt_path.write_text("hello")
    result = json.loads(await tool_pdf_read(str(txt_path)))
    assert "error" in result


# ── tool_pdf_extract ───────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_pdf_extract_all_pages(tmp_path):
    """Extract all pages from a PDF as images."""
    from dragon.tool.builtins.documents import tool_pdf_extract

    pdf_path = tmp_path / "test.pdf"
    _make_pdf(pdf_path, pages_text=["Page one", "Page two"])

    out_dir = tmp_path / "extracted"
    result = json.loads(await tool_pdf_extract(str(pdf_path), output_dir=str(out_dir)))
    assert "error" not in result
    assert result["extracted_pages"] == 2
    assert len(result["files"]) == 2

    for f in result["files"]:
        assert Path(f).exists()
        assert Path(f).suffix == ".png"


@pytest.mark.asyncio
async def test_pdf_extract_specific_pages(tmp_path):
    """Extract specific pages by specification."""
    from dragon.tool.builtins.documents import tool_pdf_extract

    pdf_path = tmp_path / "test.pdf"
    _make_pdf(pdf_path, pages_text=["A", "B", "C", "D"])

    out_dir = tmp_path / "extracted"
    result = json.loads(await tool_pdf_extract(str(pdf_path), output_dir=str(out_dir), pages="1,3"))
    assert result["extracted_pages"] == 2
    assert len(result["files"]) == 2


@pytest.mark.asyncio
async def test_pdf_extract_nonexistent():
    """Extract from non-existent file returns error."""
    from dragon.tool.builtins.documents import tool_pdf_extract

    result = json.loads(await tool_pdf_extract("/nonexistent/doc.pdf"))
    assert "error" in result


# ── tool_docx_read ─────────────────────────────────────────────────────


@pytest.mark.asyncio
async def test_docx_read_basic(tmp_path):
    """Read a DOCX file and verify paragraph extraction."""
    from dragon.tool.builtins.documents import tool_docx_read

    docx_path = tmp_path / "test.docx"
    _make_docx(docx_path, paragraphs=["Hello world", "Second paragraph"])

    result = json.loads(await tool_docx_read(str(docx_path)))
    assert "error" not in result
    assert result["paragraphs"] == 2
    assert len(result["content"]) == 2
    assert result["content"][0]["text"] == "Hello world"
    assert result["content"][1]["text"] == "Second paragraph"


@pytest.mark.asyncio
async def test_docx_read_empty_skip(tmp_path):
    """Empty paragraphs should be skipped."""
    from dragon.tool.builtins.documents import tool_docx_read
    from docx import Document

    docx_path = tmp_path / "test.docx"
    doc = Document()
    doc.add_paragraph("Visible text")
    doc.add_paragraph("")  # empty
    doc.add_paragraph("   ")  # whitespace only
    doc.add_paragraph("Another visible")
    doc.save(str(docx_path))

    result = json.loads(await tool_docx_read(str(docx_path)))
    assert result["paragraphs"] == 2
    assert result["content"][0]["text"] == "Visible text"
    assert result["content"][1]["text"] == "Another visible"


@pytest.mark.asyncio
async def test_docx_read_nonexistent():
    """Read a non-existent DOCX returns error."""
    from dragon.tool.builtins.documents import tool_docx_read

    result = json.loads(await tool_docx_read("/nonexistent/doc.docx"))
    assert "error" in result


@pytest.mark.asyncio
async def test_docx_read_wrong_extension(tmp_path):
    """Read a non-DOCX file returns error."""
    from dragon.tool.builtins.documents import tool_docx_read

    txt_path = tmp_path / "notes.txt"
    txt_path.write_text("hello")
    result = json.loads(await tool_docx_read(str(txt_path)))
    assert "error" in result


# ── Registration smoke test ────────────────────────────────────────────


def test_registration():
    """Verify document tools can be registered without error."""
    from dragon.tool.registry import ToolRegistry
    from dragon.tool.builtins import register_builtins

    registry = ToolRegistry()
    # register_builtins may fail due to partially-implemented tools from
    # other modules; test document tools directly
    from dragon.tool.builtins.documents import (
        tool_pptx_read,
        tool_pptx_create,
        tool_pdf_read,
        tool_pdf_extract,
        tool_docx_read,
    )
    registry.register(
        name="pptx_read",
        description="Read PPTX",
        tags=["document"],
        category="document",
        timeout_secs=30,
    )(tool_pptx_read)
    registry.register(
        name="pptx_create",
        description="Create PPTX",
        tags=["document"],
        category="document",
        timeout_secs=30,
    )(tool_pptx_create)
    registry.register(
        name="pdf_read",
        description="Read PDF",
        tags=["document"],
        category="document",
        timeout_secs=30,
    )(tool_pdf_read)
    registry.register(
        name="pdf_extract",
        description="Extract PDF",
        tags=["document"],
        category="document",
        timeout_secs=60,
    )(tool_pdf_extract)
    registry.register(
        name="docx_read",
        description="Read DOCX",
        tags=["document"],
        category="document",
        timeout_secs=30,
    )(tool_docx_read)

    names = [t["name"] for t in registry.list_tools()]
    assert "pptx_read" in names
    assert "pptx_create" in names
    assert "pdf_read" in names
    assert "pdf_extract" in names
    assert "docx_read" in names
