"""
Dragon Agent — Document Tools (PPT/PDF/DOCX)
=============================================

Tools for reading and creating PPTX, PDF, and DOCX documents.

Tools:
    - pptx_read: Read PPTX slides (title + text per slide)
    - pptx_create: Create a new PPTX file from slide definitions
    - pdf_read: Read text from PDF pages
    - pdf_extract: Extract pages as images from a PDF
    - docx_read: Read paragraphs from a DOCX file

Dependencies (all soft):
    - python-pptx: For PPTX read/create.
    - pymupdf (fitz): For PDF read/extract.
    - python-docx: For DOCX read.
"""
from __future__ import annotations

import json
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger("dragon.tool.builtins.documents")


# ────────────────────────────────────────────────────────────────────
# Helpers
# ────────────────────────────────────────────────────────────────────


def _resolve_path(path: str) -> Path:
    """Resolve and validate a local file path. Returns Path or raises."""
    p = Path(path).expanduser().resolve()
    if not p.exists():
        raise FileNotFoundError(str(p))
    if not p.is_file():
        raise IsADirectoryError(f"Not a file: {p}")
    return p


def _json_error(message: str, **extra: Any) -> str:
    """Return a JSON error response."""
    result: Dict[str, Any] = {"error": message}
    result.update(extra)
    return json.dumps(result, default=str)


def _json_ok(data: Dict[str, Any]) -> str:
    """Return a JSON success response."""
    return json.dumps(data, default=str)


# ────────────────────────────────────────────────────────────────────
# Tool: pptx_read
# ────────────────────────────────────────────────────────────────────


async def tool_pptx_read(path: str) -> str:
    """Read a PPTX file, returning the title and text of each slide.

    Args:
        path: Path to the .pptx file.

    Returns:
        JSON string with slide count and per-slide title + text content.
    """
    try:
        p = _resolve_path(path)
    except FileNotFoundError:
        return _json_error(f"File not found: {path}")
    except IsADirectoryError:
        return _json_error(f"Path is a directory: {path}")

    if p.suffix.lower() not in (".pptx", ".ppt"):
        return _json_error(f"Not a PPTX file: {p.suffix}")

    try:
        from pptx import Presentation
        from pptx.util import Inches  # noqa: F401
    except ImportError:
        return _json_error(
            "python-pptx is not installed",
            hint="Install with: pip install python-pptx",
        )

    try:
        prs = Presentation(str(p))
    except Exception as e:
        return _json_error(f"Failed to open PPTX: {e}")

    slides_out: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_data: Dict[str, Any] = {
            "slide": idx,
            "title": "",
            "content": [],
        }

        # Extract title (first title shape found)
        title_texts: List[str] = []
        body_texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()  # type: ignore[union-attr]
                if not text:
                    continue
                # Heuristic: shapes named "Title" or on a title placeholder are titles
                is_title = shape.is_placeholder and shape.placeholder_format.type == 1  # TITLE
                if is_title or "title" in (getattr(shape, 'name', '') or '').lower():
                    title_texts.append(text)
                else:
                    body_texts.append(text)

        if title_texts:
            slide_data["title"] = "\n".join(title_texts)
        else:
            # If no explicit title shape, use first text block as title
            if body_texts:
                slide_data["title"] = body_texts[0]
                body_texts = body_texts[1:]

        slide_data["content"] = body_texts
        slides_out.append(slide_data)

    return _json_ok({
        "file": str(p),
        "slides": len(slides_out),
        "data": slides_out,
    })


# ────────────────────────────────────────────────────────────────────
# Tool: pptx_create
# ────────────────────────────────────────────────────────────────────


async def tool_pptx_create(path: str, slides: List[Dict[str, Any]]) -> str:
    """Create a PPTX file from a list of slide definitions.

    Each slide dict should have:
        - "title": The slide title text.
        - "content": A list of strings for the body content.

    Example:
        slides = [{"title": "Hello", "content": ["Bullet 1", "Bullet 2"]}]

    Args:
        path: Output path for the .pptx file.
        slides: List of slide dicts.

    Returns:
        JSON string with success status and file path.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return _json_error(
            "python-pptx is not installed",
            hint="Install with: pip install python-pptx",
        )

    if not slides:
        return _json_error("slides list is empty", hint='Provide at least one slide dict like {"title":"...", "content":["..."]}')

    try:
        prs = Presentation()

        for i, slide_def in enumerate(slides):
            if not isinstance(slide_def, dict):
                return _json_error(f"Slide {i} is not a dict: {type(slide_def).__name__}")

            title = slide_def.get("title", "")
            content = slide_def.get("content", [])

            if isinstance(content, str):
                content = [content]
            if not isinstance(content, list):
                return _json_error(f"Slide {i + 1} content must be a list or string")

            # Use blank layout
            blank_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(blank_layout)

            # Add title textbox
            if title:
                left = Inches(1)
                top = Inches(0.5)
                width = Inches(8.5)
                height = Inches(1.2)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = str(title)
                p.font.size = Pt(32)
                p.font.bold = True

            # Add content textbox
            if content:
                left = Inches(1)
                top = Inches(2.0)
                width = Inches(8.5)
                height = Inches(4.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                for j, line in enumerate(content):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(line)
                    p.font.size = Pt(18)
                    p.space_after = Pt(8)

        # Ensure output directory exists
        out_path = Path(path).expanduser().resolve()
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))

        return _json_ok({
            "file": str(out_path),
            "slides_created": len(slides),
            "size_bytes": out_path.stat().st_size,
        })
    except Exception as e:
        logger.exception("pptx_create failed")
        return _json_error(f"Failed to create PPTX: {e}")


# ────────────────────────────────────────────────────────────────────
# Tool: pdf_read
# ────────────────────────────────────────────────────────────────────


async def tool_pdf_read(
    path: str,
    page_start: int = 1,
    page_end: Optional[int] = None,
) -> str:
    """Read text from a PDF file, optionally within a page range.

    Uses pymupdf (fitz) for high-quality text extraction.

    Args:
        path: Path to the .pdf file.
        page_start: First page to read (1-indexed, default: 1).
        page_end: Last page to read inclusive (default: all pages).

    Returns:
        JSON string with total pages and per-page text.
    """
    try:
        p = _resolve_path(path)
    except FileNotFoundError:
        return _json_error(f"File not found: {path}")
    except IsADirectoryError:
        return _json_error(f"Path is a directory: {path}")

    if p.suffix.lower() != ".pdf":
        return _json_error(f"Not a PDF file: {p.suffix}")

    try:
        import fitz
    except ImportError:
        return _json_error(
            "pymupdf (fitz) is not installed",
            hint="Install with: pip install pymupdf",
        )

    try:
        doc = fitz.open(str(p))
    except Exception as e:
        return _json_error(f"Failed to open PDF: {e}")

    try:
        total_pages = len(doc)
        if total_pages == 0:
            return _json_ok({"file": str(p), "total_pages": 0, "pages": []})

        start = max(1, page_start) - 1  # to 0-indexed
        if page_end is None:
            end = total_pages
        else:
            end = min(total_pages, page_end)

        if start >= total_pages:
            return _json_error(f"page_start {page_start} exceeds total pages {total_pages}")

        pages_out: List[Dict[str, Any]] = []
        for i in range(start, end):
            page = doc[i]
            raw_text: str = page.get_text()  # type: ignore[assignment]
            text = raw_text.strip()
            pages_out.append({
                "page": i + 1,
                "text": text,
                "chars": len(text),
            })

        return _json_ok({
            "file": str(p),
            "total_pages": total_pages,
            "page_start": start + 1,
            "page_end": end,
            "pages": pages_out,
        })
    finally:
        doc.close()


# ────────────────────────────────────────────────────────────────────
# Tool: pdf_extract
# ────────────────────────────────────────────────────────────────────


async def tool_pdf_extract(
    path: str,
    output_dir: Optional[str] = None,
    pages: Optional[str] = None,
) -> str:
    """Extract pages from a PDF as images (PNG format).

    Uses pymupdf (fitz) for rendering.

    Args:
        path: Path to the .pdf file.
        output_dir: Output directory for extracted images (default: same dir as PDF).
        pages: Page spec (e.g., "1", "1-3", "1,3,5"). None extracts all pages.

    Returns:
        JSON string with list of extracted image paths.
    """
    try:
        p = _resolve_path(path)
    except FileNotFoundError:
        return _json_error(f"File not found: {path}")
    except IsADirectoryError:
        return _json_error(f"Path is a directory: {path}")

    if p.suffix.lower() != ".pdf":
        return _json_error(f"Not a PDF file: {p.suffix}")

    try:
        import fitz
    except ImportError:
        return _json_error(
            "pymupdf (fitz) is not installed",
            hint="Install with: pip install pymupdf",
        )

    # Determine output directory
    if output_dir:
        out_dir = Path(output_dir).expanduser().resolve()
    else:
        out_dir = p.parent
    out_dir.mkdir(parents=True, exist_ok=True)

    # Parse page specification
    try:
        doc = fitz.open(str(p))
    except Exception as e:
        return _json_error(f"Failed to open PDF: {e}")

    try:
        total_pages = len(doc)

        if total_pages == 0:
            return _json_ok({"file": str(p), "total_pages": 0, "extracted": []})

        # Parse pages parameter
        page_indices: List[int] = []
        if pages is None:
            page_indices = list(range(total_pages))
        else:
            for part in pages.split(","):
                part = part.strip()
                if not part:
                    continue
                if "-" in part:
                    a, b = part.split("-", 1)
                    a, b = int(a.strip()), int(b.strip())
                    for pg in range(a, b + 1):
                        if 1 <= pg <= total_pages:
                            page_indices.append(pg - 1)
                else:
                    pg = int(part)
                    if 1 <= pg <= total_pages:
                        page_indices.append(pg - 1)

        # Remove duplicates and sort
        page_indices = sorted(set(page_indices))

        if not page_indices:
            return _json_error("No valid pages to extract", total_pages=total_pages, pages_spec=pages)

        # Extract each page as PNG
        stem = p.stem
        extracted: List[str] = []
        for idx in page_indices:
            page = doc[idx]
            # Render at a reasonable DPI
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom (~144 DPI)
            pix = page.get_pixmap(matrix=mat)
            img_path = out_dir / f"{stem}_page_{idx + 1:03d}.png"
            pix.save(str(img_path))
            extracted.append(str(img_path))

        return _json_ok({
            "file": str(p),
            "total_pages": total_pages,
            "extracted_pages": len(extracted),
            "output_dir": str(out_dir),
            "files": extracted,
        })
    finally:
        doc.close()


# ────────────────────────────────────────────────────────────────────
# Tool: docx_read
# ────────────────────────────────────────────────────────────────────


async def tool_docx_read(path: str) -> str:
    """Read text content from a DOCX file.

    Extracts all paragraphs with their style information.

    Args:
        path: Path to the .docx file.

    Returns:
        JSON string with paragraph count and content.
    """
    try:
        p = _resolve_path(path)
    except FileNotFoundError:
        return _json_error(f"File not found: {path}")
    except IsADirectoryError:
        return _json_error(f"Path is a directory: {path}")

    if p.suffix.lower() not in (".docx",):
        return _json_error(f"Not a DOCX file: {p.suffix}")

    try:
        from docx import Document
    except ImportError:
        return _json_error(
            "python-docx is not installed",
            hint="Install with: pip install python-docx",
        )

    try:
        doc = Document(str(p))
    except Exception as e:
        return _json_error(f"Failed to open DOCX: {e}")

    paragraphs_out: List[Dict[str, Any]] = []
    for idx, para in enumerate(doc.paragraphs, 1):
        text = para.text
        if not text.strip():
            continue  # Skip empty paragraphs
        style = para.style.name if para.style else "Normal"
        paragraphs_out.append({
            "index": idx,
            "style": style,
            "text": text,
        })

    # Also extract table text if any
    tables_out: List[Dict[str, Any]] = []
    for t_idx, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows.append(cells)
        tables_out.append({
            "table": t_idx,
            "rows": len(rows),
            "data": rows,
        })

    return _json_ok({
        "file": str(p),
        "paragraphs": len(paragraphs_out),
        "tables": len(tables_out),
        "content": paragraphs_out,
        "table_content": tables_out,
    })
